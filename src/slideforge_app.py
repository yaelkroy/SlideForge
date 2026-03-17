## Repository Location

### Current local working path
Current canonical local repo root at the time of writing:

`C:\temp\ML_Slides\slideforge`

This path is a working reference, not a permanent guarantee.
If the repo is moved, update this section.

### Repo identity rule
When working on this project, treat the repository root as the directory containing:
- `LLM_CONTEXT.md`
- `README.md`
- `src/`
- `docs/`
- `tests/`

from __future__ import annotations

from collections import defaultdict
from pathlib import Path
from typing import Any
import math

import numpy as np
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt




# ============================================================
# Paths
# ============================================================

BASE_DIR = Path(__file__).resolve().parent
BACKGROUND_DIR = BASE_DIR / "backgrounds"
OUTPUT_FILE = BASE_DIR / "ML_Foundations_Auto.pptx"
GENERATED_DIR = BASE_DIR / "_generated"
GENERATED_DIR.mkdir(exist_ok=True)

# ============================================================
# Slide size
# ============================================================

SLIDE_W = 13.333
SLIDE_H = 7.5


# ============================================================
# Fonts
# ============================================================

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"
FORMULA_FONT = "Cambria"


# ============================================================
# Colors
# ============================================================

NAVY = RGBColor(23, 36, 74)
SLATE = RGBColor(68, 78, 96)
MUTED = RGBColor(95, 102, 117)
ACCENT = RGBColor(120, 142, 205)

WHITE = RGBColor(255, 255, 255)
OFFWHITE = RGBColor(247, 249, 252)

BOX_LINE = RGBColor(173, 185, 220)
LIGHT_BOX_FILL = RGBColor(255, 255, 255)
DARK_BOX_FILL = RGBColor(23, 32, 56)


# ============================================================
# Background registry
# Use exact filenames already generated
# ============================================================

BACKGROUND_SETS = {
    "title": [
        "Background 8.png",
        "Background 9.png",
    ],
    "section": [
        "Background 8.png",
        "Background 9.png",
    ],
    "formula": [
        "Background 10.png",
        "Background 7.png",
        "Background 2.png",
    ],
    "concept": [
        "Background 1.png",
        "Background 4.png",
        "Background 6.png",
    ],
    "example": [
        "Background 7.png",
        "Background 1.png",
        "Background 6.png",
    ],
    "ml": [
        "Background 5.png",
        "Background 3.png",
    ],
}


# ============================================================
# Sample slide data
# Replace this list slide by slide as we build your deck
# ============================================================

SLIDES: list[dict[str, Any]] = [
{
    "kind": "title_composite",
    "theme": "title",
    "background": "Background 8.png",

    "slide_number": 1,
    "slide_title": "Machine Learning Foundations: Geometry, Optimization, and Linear Classification",
    "purpose": (
        "Open the presentation with a strong conceptual bridge: this deck is about how "
        "geometric objects, optimization ideas, and computational representations become "
        "the language of early machine learning."
    ),
    "visual": (
        "A polished composite visual with four connected layers: "
        "point/vector in space, plane or separating line, loss curve/optimization surface, "
        "and classifier boundary over labeled points. The visual should feel like a "
        "progression from math objects to ML decisions."
    ),
    "text_explanation": (
        "From vectors and planes to loss functions, feature vectors, and classifiers"
    ),
    "bullets": [
        "Geometry gives us representation",
        "Optimization gives us learning",
        "Classification gives us prediction",
    ],
    "formulas": [
        "x ∈ R^d",
        "θ · x + θ₀ = 0",
        "θ ← θ − γ∇L(θ)",
    ],
    "concrete_example_anchor": (
        "Tiny embedded labels in the visual: x=(3,2,2), 3x₁+x₂−1=0, h(x)∈{−1,+1}."
    ),
    "speaker_intent": (
        "This deck starts before the classifier. First we need the language: points, "
        "vectors, planes, loss, gradients, and features."
    ),

    # rendered text
    "title": "Machine Learning Foundations: Geometry, Optimization,\nand Linear Classification",
    "subtitle": "From vectors and planes to loss functions, feature vectors, and classifiers",
    "tiny_footer": "Foundations for early machine learning concepts and linear classification",
    "formula_ribbon": "",

    # title-slide behavior
    "show_footer_author": False,
    "show_author_line": True,
    "author_line": "Dr. Yael Demedetskaya",

    "layout": {
        "title_y": 0.90,
        "subtitle_y": 2.02,
        "author_y": 2.70,
        "visual_region": {"x": 1.00, "y": 3.30, "w": 11.30, "h": 2.00},
        "bullets_region": {"x": 2.75, "y": 5.38, "w": 7.90, "h": 0.36},
        "tiny_footer_region": {"x": 2.00, "y": 6.36, "w": 9.35, "h": 0.22},
    },

    "composite_visual": {
        "type": "connected_four_stage_banner",
        "style": "clean_academic_glass_panels",
        "panels": [
            {
                "label": "Point / Vector",
                "x": 1.05,
                "y": 3.62,
                "w": 2.30,
                "h": 1.40,
                "mini_visual": "vector_point_plane_combo",
                "embedded_label": "x=(3,2,2)",
            },
            {
                "label": "Plane / Separator",
                "x": 3.82,
                "y": 3.62,
                "w": 2.30,
                "h": 1.40,
                "mini_visual": "plane_slice_with_normal",
                "embedded_label": "3x₁+x₂−1=0",
            },
            {
                "label": "Loss / Optimization",
                "x": 6.59,
                "y": 3.62,
                "w": 2.30,
                "h": 1.40,
                "mini_visual": "loss_curve_with_descent_arrow",
                "embedded_label": "θ ← θ−γ∇L(θ)",
            },
            {
                "label": "Classifier",
                "x": 9.36,
                "y": 3.62,
                "w": 2.30,
                "h": 1.40,
                "mini_visual": "scatter_with_boundary_and_classes",
                "embedded_label": "h(x)∈{−1,+1}",
            },
        ],
        "connectors": [
            {"from": 0, "to": 1, "style": "soft_arrow"},
            {"from": 1, "to": 2, "style": "soft_arrow"},
            {"from": 2, "to": 3, "style": "soft_arrow"},
        ],
    },

    "design_notes": [
        "No footer author on title slide.",
        "No formula ribbon on title slide.",
        "Hero element is the four-panel progression.",
        "Keep lower area clean except for the three pill tags and tiny footer.",
    ],
},
{
    "kind": "section_divider",
    "theme": "section",
    "background": "Background 9.png",

    "slide_number": 2,
    "slide_title": "Part I — Framing, Notation, and Roadmap",
    "purpose": (
        "Mark the start of the first section clearly. This should feel like a section divider, "
        "not a content-heavy slide."
    ),
    "visual": (
        "Minimal, elegant section-divider design with a faint geometric mood and three soft "
        "ghost visuals: vector arrow, plane slice, scatterplot with separator."
    ),
    "text_explanation": (
        "How the course’s mathematical language becomes machine learning language"
    ),
    "bullets": [],
    "formulas": [],
    "concrete_example_anchor": (
        "Optional faint background labels: vector, hyperplane, classifier."
    ),
    "speaker_intent": (
        "In this part we establish the vocabulary and the map of the deck."
    ),

    # rendered text
    "title": "Part I — Framing, Notation, and Roadmap",
    "subtitle": "How the course’s mathematical language becomes machine learning language",
    "tiny_footer": "",
    "formula_ribbon": "",
    "show_footer_author": True,

    "layout": {
        "title_region": {"x": 1.0, "y": 2.00, "w": 11.0, "h": 0.90},
        "subtitle_region": {"x": 1.45, "y": 2.88, "w": 10.1, "h": 0.40},
        "ghost_visual_region": {"x": 1.20, "y": 3.85, "w": 10.90, "h": 1.30},
    },

    "section_visual": {
        "type": "ghost_geometry_overlay",
        "style": "minimal_dark_academic",
        "elements": [
            {
                "kind": "vector_arrow",
                "x": 1.70,
                "y": 4.00,
                "w": 2.00,
                "h": 0.80,
                "label": "vector",
            },
            {
                "kind": "plane_slice",
                "x": 5.15,
                "y": 3.92,
                "w": 2.00,
                "h": 0.84,
                "label": "hyperplane",
            },
            {
                "kind": "scatter_separator",
                "x": 8.50,
                "y": 3.92,
                "w": 2.15,
                "h": 0.84,
                "label": "classifier",
            },
        ],
        "faint_grid": True,
        "soft_connector_line": True,
    },

    "design_notes": [
        "No bullet content on the divider.",
        "The three ghost visuals should be large, simple, and readable.",
        "Labels should be tiny and decorative, not dominant.",
    ],
},
{
    "kind": "dependency_map",
    "theme": "concept",
    "background": "Background 10.png",

    "slide_number": 3,
    "slide_title": "Why These Foundations Come First",
    "purpose": (
        "Explain why geometry, calculus, probability, and computation appear before or "
        "alongside ML concepts. The source material motivates points, vectors, planes, "
        "optimization, feature vectors, classifiers, and training behavior in one shared language."
    ),
    "visual": (
        "A clean dependency map with four prerequisite areas feeding into Machine Learning, "
        "plus a compact explanatory right column."
    ),
    "text_explanation": (
        "Machine learning is not a single idea. It depends on how we represent data, "
        "measure error, optimize parameters, and reason about uncertainty."
    ),
    "bullets": [
        "Vectors represent examples and parameters",
        "Planes represent decision boundaries",
        "Loss functions define prediction error",
        "Gradients tell us how to improve",
        "Probability helps us model uncertainty",
        "Computation makes the math executable",
    ],
    "formulas": [],
    "concrete_example_anchor": (
        "movie → feature vector; line/plane → classifier boundary; "
        "bowl surface → optimization; Gaussian curve → uncertainty"
    ),
    "speaker_intent": (
        "Before a model can learn, it needs representation, error measurement, "
        "and a method for improvement."
    ),

    "title": "Why These Foundations Come First",
    "subtitle": "",
    "tiny_footer": "",
    "formula_ribbon": "",
    "show_footer_author": True,

    "layout": {
        "title_y": 0.42,

        # overall placement knobs
        "box_title_gap": 0.34,
        "box_title_h": 0.24,
        "box_title_font_size": 13,

        "box_inner_pad_x": 0.16,
        "box_inner_pad_y": 0.16,

        "note_text_font_size": 13,
        "bullet_font_size": 11,
        "bullet_sub_font_size": 10,
        "takeaway_font_size": 12,

        "connector_width_pt": 1.3,

        # left diagram block
        "diagram_region": {"x": 1.15, "y": 1.30, "w": 6.95, "h": 4.15},

        # right notes column
        "right_column_x": 8.52,
        "right_column_w": 3.58,

        "explanation_box": {"x": 8.52, "y": 1.46, "w": 3.58, "h": 0.86},
        "bullets_box": {"x": 8.40, "y": 2.88, "w": 3.70, "h": 2.38},

        "takeaway_box": {"x": 1.25, "y": 5.90, "w": 10.85, "h": 0.42},
    },

    "diagram": {
        "type": "hub_and_spokes",
        "style": "balanced_left_diagram_right_notes",

        "center_node": {
            "label": "Machine\nLearning",
            "x": 3.78,
            "y": 2.42,
            "w": 1.95,
            "h": 0.95,
            "style": "primary_hub_glow",
            "font_size": 18,
        },

        "input_nodes": [
            {
                "label": "Linear Algebra /\nGeometry",
                "x": 1.50,
                "y": 1.45,
                "w": 2.15,
                "h": 0.82,
                "style": "input_node",
                "font_size": 14,
                "mini_visual": "",
                "callout": "",
            },
            {
                "label": "Calculus /\nOptimization",
                "x": 1.50,
                "y": 3.72,
                "w": 2.15,
                "h": 0.82,
                "style": "input_node",
                "font_size": 14,
                "mini_visual": "",
                "callout": "",
            },
            {
                "label": "Probability /\nUncertainty",
                "x": 5.90,
                "y": 1.45,
                "w": 1.95,
                "h": 0.82,
                "style": "input_node",
                "font_size": 14,
                "mini_visual": "",
                "callout": "",
            },
            {
                "label": "Computation /\nNumPy",
                "x": 5.90,
                "y": 3.72,
                "w": 1.95,
                "h": 0.82,
                "style": "input_node",
                "font_size": 14,
                "mini_visual": "",
                "callout": "",
            },
        ],

        "arrows": [
            {"from": "Linear Algebra /\nGeometry", "to": "Machine\nLearning", "style": "clean_connector"},
            {"from": "Calculus /\nOptimization", "to": "Machine\nLearning", "style": "clean_connector"},
            {"from": "Probability /\nUncertainty", "to": "Machine\nLearning", "style": "clean_connector"},
            {"from": "Computation /\nNumPy", "to": "Machine\nLearning", "style": "clean_connector"},
        ],
    },

    "explanation_box": {
        "title": "Core idea",
        "text": (
            "Machine learning is built from representation, geometry, "
            "optimization, uncertainty, and computation."
        ),
    },

    "right_panel_title": "Why these prerequisites matter",
    "right_panel_bullets": [
        "Vectors represent examples and parameters",
        "Planes represent decision boundaries",
        "Loss functions define prediction error",
        "Gradients tell us how to improve",
        "Probability helps us model uncertainty",
        "Computation makes the math executable",
    ],

    "takeaway": (
        "Before a model can learn, it needs representation, error measurement, "
        "and a method for improvement."
    ),

    "design_notes": [
        "Full right column stays.",
        "Box title spacing is controlled by layout.",
        "No mini visuals or callouts on this slide.",
        "Diagram stays left; notes stay right; takeaway stays below.",
    ],
}
]


# ============================================================
# Core utilities
# ============================================================

def inches(value: float):
    return Inches(value)


def pt(value: float):
    return Pt(value)


def ensure_background_exists(filename: str) -> Path:
    path = BACKGROUND_DIR / filename
    if not path.exists():
        raise FileNotFoundError(f"Background file not found: {path}")
    return path


def validate_backgrounds() -> None:
    required = set()

    for items in BACKGROUND_SETS.values():
        for filename in items:
            required.add(filename)

    for slide in SLIDES:
        if "background" in slide:
            required.add(slide["background"])

    missing = [str(BACKGROUND_DIR / name) for name in sorted(required) if not (BACKGROUND_DIR / name).exists()]
    if missing:
        raise FileNotFoundError(
            "These background files are missing:\n- " + "\n- ".join(missing)
        )


def create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)
    return prs


def choose_background(theme: str, counters: dict[str, int]) -> str:
    options = BACKGROUND_SETS[theme]
    index = counters[theme] % len(options)
    counters[theme] += 1
    return options[index]


def add_background(slide, prs: Presentation, filename: str) -> None:
    bg_path = ensure_background_exists(filename)
    slide.shapes.add_picture(
        str(bg_path),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def add_footer(slide, text: str = "Dr. Yael Demedetskaya", dark: bool = False) -> None:
    color = OFFWHITE if dark else MUTED
    box = slide.shapes.add_textbox(inches(0.45), inches(6.95), inches(4.2), inches(0.25))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = pt(10)
    run.font.color.rgb = color


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_name: str = BODY_FONT,
    font_size: int = 18,
    color: RGBColor = NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box, tf


def style_paragraph(paragraph, font_name: str, font_size: int, color: RGBColor, bold: bool = False) -> None:
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_bullets_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: list[str | tuple[str, int]],
    color: RGBColor = NAVY,
    top_font_size: int = 20,
    sub_font_size: int = 17,
) -> None:
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        p.alignment = PP_ALIGN.LEFT
        prefix = "• " if level == 0 else "– "
        indent = "    " * level
        p.text = f"{indent}{prefix}{text}"
        p.space_after = pt(8 if level == 0 else 4)
        p.line_spacing = 1.15

        style_paragraph(
            p,
            font_name=BODY_FONT,
            font_size=top_font_size if level == 0 else sub_font_size,
            color=color,
            bold=False,
        )


def add_divider_line(slide, dark: bool = False) -> None:
    color = OFFWHITE if dark else ACCENT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(0.8), inches(1.1), inches(2.4), inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_rounded_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    line_color: RGBColor = BOX_LINE,
    fill_color: RGBColor | None = LIGHT_BOX_FILL,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(w), inches(h))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = pt(1.25)
    return shape


def add_box_title(slide, x: float, y: float, w: float, text: str, dark: bool = False) -> None:
    color = OFFWHITE if dark else SLATE
    add_textbox(
        slide,
        x=x,
        y=y,
        w=w,
        h=0.20,
        text=text,
        font_name=BODY_FONT,
        font_size=11,
        color=color,
        bold=True,
    )


# ============================================================
# Slide builders
# ============================================================

def new_slide(prs: Presentation, background_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
    add_background(slide, prs, background_name)
    return slide


def build_title_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    bg = spec.get("background") or choose_background("title", counters)
    slide = new_slide(prs, bg)

    # Main title
    add_textbox(
        slide,
        x=0.75,
        y=1.0,
        w=11.8,
        h=1.1,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Subtitle
    subtitle = spec.get("subtitle", "")
    if subtitle:
        add_textbox(
            slide,
            x=1.2,
            y=2.15,
            w=10.9,
            h=0.45,
            text=subtitle,
            font_name=BODY_FONT,
            font_size=18,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    # Author line only on title slide body
    add_textbox(
        slide,
        x=3.7,
        y=2.9,
        w=5.9,
        h=0.4,
        text="Dr. Yael Demedetskaya",
        font_name=BODY_FONT,
        font_size=18,
        color=OFFWHITE,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    # Bullets
    bullets = spec.get("bullets", [])
    if bullets:
        add_bullets_box(
            slide,
            x=2.9,
            y=3.75,
            w=7.5,
            h=1.65,
            bullets=bullets,
            color=OFFWHITE,
            top_font_size=18,
            sub_font_size=15,
        )

    # Formula ribbon
    formula_ribbon = spec.get("formula_ribbon", "")
    if formula_ribbon:
        add_textbox(
            slide,
            x=1.2,
            y=5.95,
            w=10.9,
            h=0.3,
            text=formula_ribbon,
            font_name=FORMULA_FONT,
            font_size=13,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    # Tiny footer text for title slide only
    tiny_footer = spec.get("tiny_footer", "")
    if tiny_footer:
        add_textbox(
            slide,
            x=2.0,
            y=6.45,
            w=9.3,
            h=0.22,
            text=tiny_footer,
            font_name=BODY_FONT,
            font_size=10,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    # IMPORTANT: no bottom-left footer on title slide


def build_section_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    bg = spec.get("background") or choose_background("section", counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=1.0,
        y=2.1,
        w=11.0,
        h=0.9,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=26,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    subtitle = spec.get("subtitle", "")
    if subtitle:
        add_textbox(
            slide,
            x=1.4,
            y=3.0,
            w=10.2,
            h=0.55,
            text=subtitle,
            font_name=BODY_FONT,
            font_size=15,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    add_footer(slide, dark=True)


def build_bullets_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "concept")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.45,
        w=11.6,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    add_bullets_box(
        slide,
        x=0.95,
        y=1.45,
        w=11.2,
        h=4.95,
        bullets=spec["bullets"],
        color=NAVY,
        top_font_size=21,
        sub_font_size=17,
    )

    add_footer(slide, dark=False)


def build_formula_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "formula")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.42,
        w=11.7,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    # Left concept box
    add_box_title(slide, 0.95, 1.22, 3.25, spec.get("intro_title", "Idea"))
    add_rounded_box(slide, 0.85, 1.45, 3.35, 3.65)
    add_textbox(
        slide,
        x=1.05,
        y=1.72,
        w=2.95,
        h=3.0,
        text=spec.get("intro_text", ""),
        font_name=BODY_FONT,
        font_size=18,
        color=SLATE,
        bold=False,
    )

    # Right formula box
    add_box_title(slide, 4.55, 1.22, 7.05, spec.get("formula_title", "Formula"))
    add_rounded_box(slide, 4.45, 1.45, 7.1, 3.65)

    formula_text = "\n\n".join(spec.get("formulas", []))
    add_textbox(
        slide,
        x=4.8,
        y=1.75,
        w=6.4,
        h=2.85,
        text=formula_text,
        font_name=FORMULA_FONT,
        font_size=22,
        color=NAVY,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    # Takeaway box
    add_box_title(slide, 0.95, 5.25, 11.2, "Takeaway")
    add_rounded_box(slide, 0.85, 5.5, 11.65, 1.0)
    add_textbox(
        slide,
        x=1.08,
        y=5.78,
        w=11.1,
        h=0.45,
        text=spec.get("takeaway", ""),
        font_name=BODY_FONT,
        font_size=17,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, dark=False)


def build_two_column_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "concept")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.42,
        w=11.7,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    # Left column
    add_box_title(slide, 0.95, 1.2, 5.1, spec.get("left_title", "Left"))
    add_rounded_box(slide, 0.85, 1.45, 5.2, 4.8)
    add_bullets_box(
        slide,
        x=1.05,
        y=1.75,
        w=4.8,
        h=4.2,
        bullets=spec.get("left_items", []),
        color=NAVY,
        top_font_size=19,
        sub_font_size=16,
    )

    # Right column
    add_box_title(slide, 6.35, 1.2, 5.95, spec.get("right_title", "Right"))
    add_rounded_box(slide, 6.25, 1.45, 6.15, 4.8)
    add_textbox(
        slide,
        x=6.55,
        y=1.75,
        w=5.6,
        h=4.15,
        text=spec.get("right_text", ""),
        font_name=FORMULA_FONT,
        font_size=19,
        color=NAVY,
        bold=False,
        align=PP_ALIGN.LEFT,
    )

    add_footer(slide, dark=False)


def build_image_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "example")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.42,
        w=11.7,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    image_path = Path(spec["image_path"])
    if not image_path.is_absolute():
        image_path = BASE_DIR / image_path

    if not image_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    slide.shapes.add_picture(str(image_path), inches(1.0), inches(1.45), width=inches(11.2), height=inches(4.9))

    caption = spec.get("caption", "")
    if caption:
        add_textbox(
            slide,
            x=1.1,
            y=6.45,
            w=11.0,
            h=0.35,
            text=caption,
            font_name=BODY_FONT,
            font_size=14,
            color=SLATE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    add_footer(slide, dark=False)
def add_node_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill_color: RGBColor = LIGHT_BOX_FILL,
    line_color: RGBColor = BOX_LINE,
    text_color: RGBColor = NAVY,
    font_size: int = 18,
    bold: bool = True,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.25)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color

    return shape
def add_hub_box(slide, x: float, y: float, w: float, h: float, text: str):
    return add_node_box(
        slide,
        x=x, y=y, w=w, h=h,
        text=text,
        fill_color=RGBColor(232, 238, 252),
        line_color=ACCENT,
        text_color=NAVY,
        font_size=20,
        bold=True,
    )

def add_arrow(
    slide,
    x1: float, y1: float,
    x2: float, y2: float,
    color: RGBColor = ACCENT,
    width_pt: float = 1.75,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inches(x1), inches(y1),
        inches(x2), inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)

    # Arrowhead
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass

    return line
def add_callout(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_size: int = 11,
    color: RGBColor = SLATE,
):
    add_textbox(
        slide,
        x=x, y=y, w=w, h=h,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=color,
        bold=False,
        align=PP_ALIGN.CENTER,
    )
def add_image(
    slide,
    image_path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
):
    slide.shapes.add_picture(
        str(image_path),
        inches(x), inches(y),
        width=inches(w),
        height=inches(h),
    )
def create_gaussian_plot(filename: str = "gaussian_curve.png") -> Path:
    path = GENERATED_DIR / filename

    x = np.linspace(-4, 4, 400)
    y = (1 / np.sqrt(2 * np.pi)) * np.exp(-(x ** 2) / 2)

    fig, ax = plt.subplots(figsize=(4, 2.2))
    ax.plot(x, y)
    ax.set_xticks([])
    ax.set_yticks([])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    plt.close(fig)

    return path   
def create_loss_curve_plot(filename: str = "loss_curve.png") -> Path:
    path = GENERATED_DIR / filename

    x = np.linspace(-2.5, 2.5, 400)
    y = 0.5 * x**2 + 0.2 * x + 0.6

    fig, ax = plt.subplots(figsize=(4, 2.2))
    ax.plot(x, y)
    ax.scatter([0], [0.6], s=30)
    ax.set_xticks([])
    ax.set_yticks([])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    plt.close(fig)

    return path
def create_scatter_separator_plot(filename: str = "scatter_separator_plot.png") -> Path:
    path = GENERATED_DIR / filename

    pos = np.array([[2.2, 4.1], [3.2, 3.7], [4.0, 3.2]])
    neg = np.array([[6.0, 2.1], [6.9, 2.4], [7.8, 2.8]])

    fig, ax = plt.subplots(figsize=(3.2, 1.8))
    ax.scatter(pos[:, 0], pos[:, 1], marker="o", s=34)
    ax.scatter(neg[:, 0], neg[:, 1], marker="x", s=34)

    xx = np.linspace(1.5, 8.5, 100)
    yy = -0.35 * xx + 5.0
    ax.plot(xx, yy, linewidth=1.8)

    ax.set_xlim(1.0, 9.0)
    ax.set_ylim(1.2, 4.8)
    ax.set_xticks([])
    ax.set_yticks([])
    for side in ax.spines.values():
        side.set_visible(False)

    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return path

def add_mini_visual(slide, kind: str, x: float, y: float, w: float, h: float):
    if kind == "gaussian_curve":
        img = create_gaussian_plot("mini_gaussian.png")
        add_image(slide, img, x, y, w, h)

    elif kind == "bowl_surface_with_downhill_arrow":
        img = create_loss_curve_plot("mini_loss_curve.png")
        add_image(slide, img, x, y, w, h)

    elif kind == "vector_arrow_and_plane_slice":
        # simple native shapes
        add_arrow(slide, x + 0.15, y + h - 0.15, x + w - 0.2, y + 0.2, color=ACCENT)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            inches(x + 0.15), inches(y + h * 0.65),
            inches(w - 0.3), inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = SLATE
        line.line.color.rgb = SLATE

    elif kind == "array_or_code_brackets":
        add_textbox(
            slide,
            x=x, y=y, w=w, h=h,
            text="[1, 0, 1]\n[0, 1, 1]",
            font_name=FORMULA_FONT,
            font_size=12,
            color=SLATE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    elif kind == "scatterplot_with_separator":
        img = create_scatter_separator_plot("mini_scatter_separator.png")
        add_image(slide, img, x, y, w, h)
def build_dependency_map_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "concept")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    layout = spec.get("layout", {})
    title_y = layout.get("title_y", 0.42)

    # configurable spacing / typography knobs
    box_title_gap = layout.get("box_title_gap", 0.36)
    box_title_h = layout.get("box_title_h", 0.24)
    box_title_font_size = layout.get("box_title_font_size", 13)

    box_inner_pad_x = layout.get("box_inner_pad_x", 0.16)
    box_inner_pad_y = layout.get("box_inner_pad_y", 0.16)

    note_text_font_size = layout.get("note_text_font_size", 13)
    bullet_font_size = layout.get("bullet_font_size", 11)
    bullet_sub_font_size = layout.get("bullet_sub_font_size", 10)
    takeaway_font_size = layout.get("takeaway_font_size", 12)

    connector_width_pt = layout.get("connector_width_pt", 1.3)

    # ------------------------------------------------------------
    # slide title
    # ------------------------------------------------------------
    add_textbox(
        slide,
        x=0.8,
        y=title_y,
        w=11.7,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    # ------------------------------------------------------------
    # dependency diagram
    # ------------------------------------------------------------
    diagram = spec["diagram"]
    center = diagram["center_node"]
    nodes = diagram["input_nodes"]

    # connectors first, behind boxes
    for node in nodes:
        if node["x"] < center["x"]:
            x1 = node["x"] + node["w"]
            y1 = node["y"] + node["h"] / 2
            x2 = center["x"]
            y2 = center["y"] + center["h"] / 2
        else:
            x1 = node["x"]
            y1 = node["y"] + node["h"] / 2
            x2 = center["x"] + center["w"]
            y2 = center["y"] + center["h"] / 2

        add_soft_connector(
            slide,
            x1=x1,
            y1=y1,
            x2=x2,
            y2=y2,
            color=ACCENT,
            width_pt=connector_width_pt,
        )

    # hub
    add_hub_box(
        slide,
        x=center["x"],
        y=center["y"],
        w=center["w"],
        h=center["h"],
        text=center["label"],
    )

    # prerequisite nodes
    for node in nodes:
        add_node_box(
            slide,
            x=node["x"],
            y=node["y"],
            w=node["w"],
            h=node["h"],
            text=node["label"],
            fill_color=LIGHT_BOX_FILL,
            line_color=BOX_LINE,
            text_color=NAVY,
            font_size=node.get("font_size", 14),
            bold=True,
        )

        if node.get("mini_visual", "").strip():
            add_mini_visual(
                slide,
                kind=node["mini_visual"],
                x=node["x"] + 0.18,
                y=node["y"] - 0.58,
                w=node["w"] - 0.36,
                h=0.42,
                suffix=f"_dep_{node['label'].replace(chr(10), '_').replace('/', '_')}",
            )

        if node.get("callout", "").strip():
            add_callout(
                slide,
                x=node["x"] - 0.02,
                y=node["y"] + node["h"] + 0.03,
                w=node["w"] + 0.04,
                h=0.24,
                text=node["callout"],
                font_size=9,
                color=SLATE,
            )

    # ------------------------------------------------------------
    # explanation box
    # ------------------------------------------------------------
    exp = spec.get("explanation_box", {})
    exp_box = layout.get("explanation_box")
    if exp_box and exp and exp.get("text", "").strip():
        add_textbox(
            slide,
            x=exp_box["x"] + 0.03,
            y=exp_box["y"] - box_title_gap,
            w=exp_box["w"] - 0.06,
            h=box_title_h,
            text=exp.get("title", "Core idea"),
            font_name=BODY_FONT,
            font_size=box_title_font_size,
            color=SLATE,
            bold=True,
        )

        add_rounded_box(
            slide,
            exp_box["x"],
            exp_box["y"],
            exp_box["w"],
            exp_box["h"],
        )

        add_textbox(
            slide,
            x=exp_box["x"] + box_inner_pad_x,
            y=exp_box["y"] + box_inner_pad_y,
            w=exp_box["w"] - 2 * box_inner_pad_x,
            h=exp_box["h"] - 2 * box_inner_pad_y,
            text=exp["text"],
            font_name=BODY_FONT,
            font_size=note_text_font_size,
            color=SLATE,
            bold=False,
        )

    # ------------------------------------------------------------
    # bullets box
    # ------------------------------------------------------------
    right_bullets = spec.get("right_panel_bullets", [])
    bullets_box = layout.get("bullets_box")
    if bullets_box and right_bullets:
        add_textbox(
            slide,
            x=bullets_box["x"] + 0.03,
            y=bullets_box["y"] - box_title_gap,
            w=bullets_box["w"] - 0.06,
            h=box_title_h,
            text=spec.get("right_panel_title", "Why this matters"),
            font_name=BODY_FONT,
            font_size=box_title_font_size,
            color=SLATE,
            bold=True,
        )

        add_rounded_box(
            slide,
            bullets_box["x"],
            bullets_box["y"],
            bullets_box["w"],
            bullets_box["h"],
        )

        add_bullets_box(
            slide,
            x=bullets_box["x"] + box_inner_pad_x,
            y=bullets_box["y"] + box_inner_pad_y,
            w=bullets_box["w"] - 2 * box_inner_pad_x,
            h=bullets_box["h"] - 2 * box_inner_pad_y,
            bullets=right_bullets,
            color=NAVY,
            top_font_size=bullet_font_size,
            sub_font_size=bullet_sub_font_size,
        )

    # ------------------------------------------------------------
    # takeaway strip
    # ------------------------------------------------------------
    takeaway = spec.get("takeaway", "").strip()
    takeaway_box = layout.get("takeaway_box")
    if takeaway and takeaway_box:
        add_rounded_box(
            slide,
            takeaway_box["x"],
            takeaway_box["y"],
            takeaway_box["w"],
            takeaway_box["h"],
        )

        add_textbox(
            slide,
            x=takeaway_box["x"] + 0.16,
            y=takeaway_box["y"] + 0.06,
            w=takeaway_box["w"] - 0.32,
            h=takeaway_box["h"] - 0.08,
            text=takeaway,
            font_name=BODY_FONT,
            font_size=takeaway_font_size,
            color=SLATE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    add_footer(slide, dark=False)


# ============================================================
# Visual helper functions for composite/title/divider slides
# ============================================================

TITLE_PANEL_FILL = RGBColor(34, 46, 84)
TITLE_PANEL_LINE = RGBColor(198, 208, 236)
GHOST_TEXT = RGBColor(224, 230, 244)


def add_image(slide, image_path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(
        str(image_path),
        inches(x), inches(y),
        width=inches(w),
        height=inches(h),
    )


def add_centered_lines_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[str],
    color: RGBColor = OFFWHITE,
    font_size: int = 16,
    bold: bool = False,
) -> None:
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER
        p.text = line
        p.space_after = pt(4)
        style_paragraph(
            p,
            font_name=BODY_FONT,
            font_size=font_size,
            color=color,
            bold=bold,
        )


def add_title_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    embedded_label: str = "",
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_PANEL_FILL
    shape.line.color.rgb = TITLE_PANEL_LINE
    shape.line.width = Pt(1.2)

    add_textbox(
        slide,
        x=x + 0.08,
        y=y + 0.08,
        w=w - 0.16,
        h=0.22,
        text=title,
        font_name=BODY_FONT,
        font_size=11,
        color=OFFWHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    if embedded_label:
        add_textbox(
            slide,
            x=x + 0.08,
            y=y + h - 0.32,
            w=w - 0.16,
            h=0.18,
            text=embedded_label,
            font_name=FORMULA_FONT,
            font_size=10,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )


def add_soft_connector(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor = TITLE_PANEL_LINE,
    width_pt: float = 1.5,
) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inches(x1), inches(y1),
        inches(x2), inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)

def add_pill_tag(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill_color: RGBColor = TITLE_PANEL_FILL,
    line_color: RGBColor = TITLE_PANEL_LINE,
    text_color: RGBColor = OFFWHITE,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.0)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = pt(12)
    run.font.bold = False
    run.font.color.rgb = text_color
def add_ghost_label(
    slide,
    x: float,
    y: float,
    w: float,
    text: str,
    font_size: int = 10,
) -> None:
    add_textbox(
        slide,
        x=x,
        y=y,
        w=w,
        h=0.2,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=GHOST_TEXT,
        bold=False,
        align=PP_ALIGN.CENTER,
    )
def build_title_composite_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    bg = spec.get("background") or choose_background("title", counters)
    slide = new_slide(prs, bg)

    layout = spec.get("layout", {})

    # Title
    add_textbox(
        slide,
        x=0.75,
        y=0.92,
        w=11.8,
        h=1.0,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Subtitle
    subtitle = spec.get("subtitle", "")
    if subtitle:
        add_textbox(
            slide,
            x=1.15,
            y=2.02,
            w=11.0,
            h=0.4,
            text=subtitle,
            font_name=BODY_FONT,
            font_size=17,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    # Author
    if spec.get("show_author_line", True):
        add_textbox(
            slide,
            x=3.6,
            y=2.72,
            w=6.1,
            h=0.35,
            text=spec.get("author_line", "Dr. Yael Demedetskaya"),
            font_name=BODY_FONT,
            font_size=17,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    # Composite visual
    composite = spec.get("composite_visual", {})
    panels = composite.get("panels", [])

    for idx, panel in enumerate(panels):
        add_title_panel(
            slide,
            x=panel["x"],
            y=panel["y"],
            w=panel["w"],
            h=panel["h"],
            title=panel["label"],
            embedded_label=panel.get("embedded_label", ""),
        )

        add_mini_visual(
            slide,
            kind=panel.get("mini_visual", ""),
            x=panel["x"] + 0.22,
            y=panel["y"] + 0.28,
            w=panel["w"] - 0.44,
            h=0.66,
            suffix=f"_title_{idx}",
        )

    # connectors between panels
    connectors = composite.get("connectors", [])
    for conn in connectors:
        p1 = panels[conn["from"]]
        p2 = panels[conn["to"]]
        add_soft_connector(
            slide,
            x1=p1["x"] + p1["w"],
            y1=p1["y"] + p1["h"] / 2,
            x2=p2["x"],
            y2=p2["y"] + p2["h"] / 2,
            color=TITLE_PANEL_LINE,
            width_pt=1.4,
        )

    # Replace bullet stack with pill tags
    bullets = spec.get("bullets", [])
    if len(bullets) >= 3:
        tag_y = 5.42
        add_pill_tag(slide, 3.0, tag_y, 2.35, 0.34, bullets[0])
        add_pill_tag(slide, 5.52, tag_y, 2.28, 0.34, bullets[1])
        add_pill_tag(slide, 7.95, tag_y, 2.45, 0.34, bullets[2])

    # Tiny footer only
    tiny_footer = spec.get("tiny_footer", "")
    if tiny_footer:
        add_textbox(
            slide,
            x=2.0,
            y=6.36,
            w=9.35,
            h=0.22,
            text=tiny_footer,
            font_name=BODY_FONT,
            font_size=10,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    # IMPORTANT: no formula ribbon on title slide
    # IMPORTANT: no footer author on title slide
def build_section_divider_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    bg = spec.get("background") or choose_background("section", counters)
    slide = new_slide(prs, bg)

    # Title
    add_textbox(
        slide,
        x=1.0,
        y=2.0,
        w=11.0,
        h=0.9,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=26,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Subtitle
    subtitle = spec.get("subtitle", "")
    if subtitle:
        add_textbox(
            slide,
            x=1.45,
            y=2.9,
            w=10.1,
            h=0.4,
            text=subtitle,
            font_name=BODY_FONT,
            font_size=15,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    # three ghost visuals, larger and simpler
    elements = spec.get("section_visual", {}).get("elements", [])

    if len(elements) >= 3:
        left, middle, right = elements[:3]

        add_mini_visual(
            slide, left["kind"],
            x=1.75, y=4.0, w=2.0, h=0.82,
            suffix="_section_left"
        )
        add_mini_visual(
            slide, middle["kind"],
            x=5.15, y=3.92, w=2.0, h=0.88,
            suffix="_section_mid"
        )
        add_mini_visual(
            slide, right["kind"],
            x=8.45, y=3.92, w=2.15, h=0.88,
            suffix="_section_right"
        )

        # subtle labels above, not below
        add_ghost_label(slide, 2.05, 4.86, 1.4, "vector", font_size=10)
        add_ghost_label(slide, 5.45, 4.86, 1.4, "hyperplane", font_size=10)
        add_ghost_label(slide, 8.78, 4.86, 1.5, "classifier", font_size=10)

        # single subtle baseline connector
        add_soft_connector(
            slide,
            x1=3.95, y1=4.36,
            x2=8.35, y2=4.36,
            color=GHOST_TEXT,
            width_pt=1.0,
        )

    add_footer(slide, dark=True)

# ============================================================
# Mini plot/image generators
# ============================================================

def create_vector_point_plot(filename: str = "vector_point_plot.png") -> Path:
    path = GENERATED_DIR / filename

    fig, ax = plt.subplots(figsize=(3.2, 1.8))
    ax.set_xlim(0, 4)
    ax.set_ylim(0, 3)
    ax.arrow(0.4, 0.4, 2.2, 1.4, head_width=0.12, head_length=0.16, length_includes_head=True)
    ax.scatter([2.6], [1.8], s=20)
    ax.plot([0.4, 2.6], [0.4, 1.8], linewidth=1.2)
    ax.set_xticks([])
    ax.set_yticks([])
    for side in ax.spines.values():
        side.set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return path


def create_plane_normal_plot(filename: str = "plane_normal_plot.png") -> Path:
    path = GENERATED_DIR / filename

    fig, ax = plt.subplots(figsize=(3.2, 1.8))
    ax.set_xlim(0, 4)
    ax.set_ylim(0, 3)
    ax.plot([0.6, 3.3], [0.9, 2.2], linewidth=1.6)
    ax.arrow(2.0, 1.55, -0.45, 0.75, head_width=0.12, head_length=0.15, length_includes_head=True)
    ax.set_xticks([])
    ax.set_yticks([])
    for side in ax.spines.values():
        side.set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return path


def create_loss_curve_plot(filename: str = "loss_curve_plot.png") -> Path:
    path = GENERATED_DIR / filename

    x = np.linspace(-2.0, 2.2, 400)
    y = 0.35 * (x - 0.4) ** 2 + 0.45

    fig, ax = plt.subplots(figsize=(3.2, 1.8))
    ax.plot(x, y, linewidth=1.6)
    ax.scatter([1.3], [0.35 * (1.3 - 0.4) ** 2 + 0.45], s=20)
    ax.annotate("", xy=(0.7, 0.55), xytext=(1.15, 0.85), arrowprops=dict(arrowstyle="->", lw=1.1))
    ax.set_xticks([])
    ax.set_yticks([])
    for side in ax.spines.values():
        side.set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return path


def create_scatter_separator_plot(filename: str = "scatter_separator_plot.png") -> Path:
    path = GENERATED_DIR / filename

    pos = np.array([[1.0, 2.1], [1.6, 2.5], [2.0, 2.0]])
    neg = np.array([[3.1, 0.9], [3.4, 1.4], [2.7, 1.0]])

    fig, ax = plt.subplots(figsize=(3.2, 1.8))
    ax.scatter(pos[:, 0], pos[:, 1], marker="o", s=28)
    ax.scatter(neg[:, 0], neg[:, 1], marker="x", s=28)
    xx = np.linspace(0.6, 3.8, 100)
    yy = -0.7 * xx + 3.45
    ax.plot(xx, yy, linewidth=1.4)
    ax.set_xlim(0.4, 4.0)
    ax.set_ylim(0.4, 3.0)
    ax.set_xticks([])
    ax.set_yticks([])
    for side in ax.spines.values():
        side.set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return path


def create_array_icon_plot(filename: str = "array_icon_plot.png") -> Path:
    path = GENERATED_DIR / filename

    fig, ax = plt.subplots(figsize=(3.2, 1.8))
    ax.axis("off")
    ax.text(
        0.5, 0.5,
        "[1, 0, 1]\n[0, 1, 1]",
        ha="center", va="center",
        fontsize=14,
        family="DejaVu Sans Mono",
    )
    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return path


def add_mini_visual(slide, kind: str, x: float, y: float, w: float, h: float, suffix: str = "") -> None:
    if kind == "vector_point_plane_combo" or kind == "vector_arrow":
        img = create_vector_point_plot(f"vector_point{suffix}.png")
        add_image(slide, img, x, y, w, h)
    elif kind == "plane_slice_with_normal" or kind == "plane_slice":
        img = create_plane_normal_plot(f"plane_normal{suffix}.png")
        add_image(slide, img, x, y, w, h)
    elif kind == "loss_curve_with_descent_arrow" or kind == "bowl_surface_with_downhill_arrow":
        img = create_loss_curve_plot(f"loss_curve{suffix}.png")
        add_image(slide, img, x, y, w, h)
    elif kind == "scatter_with_boundary_and_classes" or kind == "scatter_separator":
        img = create_scatter_separator_plot(f"scatter_separator{suffix}.png")
        add_image(slide, img, x, y, w, h)
    elif kind == "array_or_code_brackets":
        img = create_array_icon_plot(f"array_icon{suffix}.png")
        add_image(slide, img, x, y, w, h)
# ============================================================
# Main build loop
# ============================================================

def build_deck() -> None:
    validate_backgrounds()

    prs = create_presentation()
    counters: dict[str, int] = defaultdict(int)


    builders = {
        "title": build_title_slide,
        "title_composite": build_title_composite_slide,
        "section": build_section_slide,
        "section_divider": build_section_divider_slide,
        "bullets": build_bullets_slide,
        "formula": build_formula_slide,
        "two_column": build_two_column_slide,
        "image": build_image_slide,
        "dependency_map": build_dependency_map_slide,
    }

    for idx, spec in enumerate(SLIDES, start=1):
        kind = spec["kind"]
        if kind not in builders:
            raise ValueError(f"Unknown slide kind on slide {idx}: {kind}")
        builders[kind](prs, spec, counters)

    prs.save(OUTPUT_FILE)
    print(f"Saved presentation: {OUTPUT_FILE}")


if __name__ == "__main__":
    build_deck()