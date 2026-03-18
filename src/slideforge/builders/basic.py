from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from slideforge.config.paths import BASE_DIR
from slideforge.config.constants import (
    TITLE_FONT,
    BODY_FONT,
    FORMULA_FONT,
    NAVY,
    SLATE,
    WHITE,
    OFFWHITE,
)
from slideforge.io.backgrounds import choose_background
from slideforge.render.primitives import (
    add_background,
    add_footer,
    add_textbox,
    add_bullets_box,
    add_divider_line,
    add_rounded_box,
    add_box_title,
)


def new_slide(prs: Presentation, background_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, background_name)
    return slide


def build_title_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    bg = spec.get("background") or choose_background("title", counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.75,
        y=1.0,
        w=11.8,
        h=1.1,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    subtitle = spec.get("subtitle", "")
    if subtitle:
        add_textbox(
            slide,
            x=1.2,
            y=2.15,
            w=10.9,
            h=0.45,
            text=subtitle,
            font_name=BODY_FONT,
            font_size=18,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        x=3.7,
        y=2.9,
        w=5.9,
        h=0.4,
        text=spec.get("author_line", "Dr. Yael Demedetskaya"),
        font_name=BODY_FONT,
        font_size=18,
        color=OFFWHITE,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    bullets = spec.get("bullets", [])
    if bullets:
        add_bullets_box(
            slide,
            x=2.9,
            y=3.75,
            w=7.5,
            h=1.65,
            bullets=bullets,
            color=OFFWHITE,
            top_font_size=18,
            sub_font_size=15,
        )

    formula_ribbon = spec.get("formula_ribbon", "")
    if formula_ribbon:
        add_textbox(
            slide,
            x=1.2,
            y=5.95,
            w=10.9,
            h=0.3,
            text=formula_ribbon,
            font_name=FORMULA_FONT,
            font_size=13,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    tiny_footer = spec.get("tiny_footer", "")
    if tiny_footer:
        add_textbox(
            slide,
            x=2.0,
            y=6.45,
            w=9.3,
            h=0.22,
            text=tiny_footer,
            font_name=BODY_FONT,
            font_size=10,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )


def build_section_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    bg = spec.get("background") or choose_background("section", counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=1.0,
        y=2.1,
        w=11.0,
        h=0.9,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=26,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    subtitle = spec.get("subtitle", "")
    if subtitle:
        add_textbox(
            slide,
            x=1.4,
            y=3.0,
            w=10.2,
            h=0.55,
            text=subtitle,
            font_name=BODY_FONT,
            font_size=15,
            color=OFFWHITE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    add_footer(slide, dark=True)


def build_bullets_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "concept")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.45,
        w=11.6,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    add_bullets_box(
        slide,
        x=0.95,
        y=1.45,
        w=11.2,
        h=4.95,
        bullets=spec["bullets"],
        color=NAVY,
        top_font_size=21,
        sub_font_size=17,
    )

    add_footer(slide, dark=False)


def build_formula_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "formula")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.42,
        w=11.7,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    add_box_title(slide, 0.95, 1.22, 3.25, spec.get("intro_title", "Idea"))
    add_rounded_box(slide, 0.85, 1.45, 3.35, 3.65)
    add_textbox(
        slide,
        x=1.05,
        y=1.72,
        w=2.95,
        h=3.0,
        text=spec.get("intro_text", ""),
        font_name=BODY_FONT,
        font_size=18,
        color=SLATE,
        bold=False,
    )

    add_box_title(slide, 4.55, 1.22, 7.05, spec.get("formula_title", "Formula"))
    add_rounded_box(slide, 4.45, 1.45, 7.1, 3.65)

    formula_text = "\n\n".join(spec.get("formulas", []))
    add_textbox(
        slide,
        x=4.8,
        y=1.75,
        w=6.4,
        h=2.85,
        text=formula_text,
        font_name=FORMULA_FONT,
        font_size=22,
        color=NAVY,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    add_box_title(slide, 0.95, 5.25, 11.2, "Takeaway")
    add_rounded_box(slide, 0.85, 5.5, 11.65, 1.0)
    add_textbox(
        slide,
        x=1.08,
        y=5.78,
        w=11.1,
        h=0.45,
        text=spec.get("takeaway", ""),
        font_name=BODY_FONT,
        font_size=17,
        color=SLATE,
        bold=False,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, dark=False)


def build_two_column_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    theme = spec.get("theme", "concept")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.42,
        w=11.7,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    add_box_title(slide, 0.95, 1.2, 5.1, spec.get("left_title", "Left"))
    add_rounded_box(slide, 0.85, 1.45, 5.2, 4.8)
    add_bullets_box(
        slide,
        x=1.05,
        y=1.75,
        w=4.8,
        h=4.2,
        bullets=spec.get("left_items", []),
        color=NAVY,
        top_font_size=19,
        sub_font_size=16,
    )

    add_box_title(slide, 6.35, 1.2, 5.95, spec.get("right_title", "Right"))
    add_rounded_box(slide, 6.25, 1.45, 6.15, 4.8)
    add_textbox(
        slide,
        x=6.55,
        y=1.75,
        w=5.6,
        h=4.15,
        text=spec.get("right_text", ""),
        font_name=FORMULA_FONT,
        font_size=19,
        color=NAVY,
        bold=False,
        align=PP_ALIGN.LEFT,
    )

    add_footer(slide, dark=False)


def build_image_slide(prs: Presentation, spec: dict[str, Any], counters: dict[str, int]) -> None:
    from slideforge.utils.units import inches

    theme = spec.get("theme", "example")
    bg = spec.get("background") or choose_background(theme, counters)
    slide = new_slide(prs, bg)

    add_textbox(
        slide,
        x=0.8,
        y=0.42,
        w=11.7,
        h=0.5,
        text=spec["title"],
        font_name=TITLE_FONT,
        font_size=24,
        color=NAVY,
        bold=True,
    )
    add_divider_line(slide, dark=False)

    image_path = Path(spec["image_path"])
    if not image_path.is_absolute():
        image_path = BASE_DIR / image_path

    if not image_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    slide.shapes.add_picture(str(image_path), inches(1.0), inches(1.45), width=inches(11.2), height=inches(4.9))

    caption = spec.get("caption", "")
    if caption:
        add_textbox(
            slide,
            x=1.1,
            y=6.45,
            w=11.0,
            h=0.35,
            text=caption,
            font_name=BODY_FONT,
            font_size=14,
            color=SLATE,
            bold=False,
            align=PP_ALIGN.CENTER,
        )

    add_footer(slide, dark=False)