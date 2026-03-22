from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from slideforge.config.constants import (
    ACCENT,
    BODY_FONT,
    BOX_LINE,
    FORMULA_FONT,
    LIGHT_BOX_FILL,
    MUTED,
    NAVY,
    OFFWHITE,
    SLATE,
    TITLE_PANEL_FILL,
    TITLE_PANEL_LINE,
)
from slideforge.utils.units import inches, pt
from slideforge.render.text import _add_text_to_shape, add_textbox



def _set_shape_line_and_fill(
    shape,
    *,
    line_color: RGBColor,
    fill_color: RGBColor | None,
    line_width_pt: float,
) -> None:
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = pt(line_width_pt)



def _add_shape(
    slide,
    *,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    line_color: RGBColor,
    fill_color: RGBColor | None,
    line_width_pt: float,
):
    shape = slide.shapes.add_shape(shape_type, inches(x), inches(y), inches(w), inches(h))
    _set_shape_line_and_fill(
        shape,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )
    return shape



def add_rounded_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    line_color: RGBColor = BOX_LINE,
    fill_color: RGBColor | None = LIGHT_BOX_FILL,
    line_width_pt: float = 1.25,
):
    return _add_shape(
        slide,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )



def add_surface_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    line_color: RGBColor = BOX_LINE,
    fill_color: RGBColor | None = LIGHT_BOX_FILL,
    line_width_pt: float = 1.25,
):
    return add_rounded_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )



def add_box_title(
    slide,
    x: float,
    y: float,
    w: float,
    text: str,
    dark: bool = False,
    *,
    color: RGBColor | None = None,
    font_size: int = 11,
    bold: bool = True,
    align=PP_ALIGN.LEFT,
) -> None:
    actual_color = color or (OFFWHITE if dark else SLATE)
    add_textbox(
        slide,
        x=x,
        y=y,
        w=w,
        h=0.20,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=actual_color,
        bold=bold,
        align=align,
    )



def add_footer(
    slide,
    text: str = "Dr. Yael Demedetskaya",
    dark: bool = False,
    *,
    x: float = 0.45,
    y: float = 6.95,
    w: float = 4.2,
    h: float = 0.25,
    color: RGBColor | None = None,
    font_size: int = 10,
    align=PP_ALIGN.LEFT,
) -> None:
    actual_color = color or (OFFWHITE if dark else MUTED)
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    _add_text_to_shape(
        box,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=actual_color,
        bold=False,
        align=align,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )



def add_node_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill_color: RGBColor = LIGHT_BOX_FILL,
    line_color: RGBColor = BOX_LINE,
    text_color: RGBColor = NAVY,
    font_size: int = 18,
    bold: bool = True,
    line_width_pt: float = 1.25,
):
    shape = add_rounded_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )
    _add_text_to_shape(
        shape,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=text_color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    return shape



def add_hub_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    fill_color: RGBColor = RGBColor(232, 238, 252),
    line_color: RGBColor = ACCENT,
    text_color: RGBColor = NAVY,
    font_size: int = 20,
    bold: bool = True,
):
    return add_node_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        text=text,
        fill_color=fill_color,
        line_color=line_color,
        text_color=text_color,
        font_size=font_size,
        bold=bold,
    )



def add_title_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    embedded_label: str = "",
    *,
    fill_color: RGBColor = TITLE_PANEL_FILL,
    line_color: RGBColor = TITLE_PANEL_LINE,
    title_color: RGBColor = OFFWHITE,
    embedded_label_color: RGBColor = OFFWHITE,
    title_font_size: int = 11,
    embedded_label_font_size: int = 10,
) -> None:
    add_rounded_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=1.2,
    )
    add_textbox(
        slide,
        x=x + 0.08,
        y=y + 0.08,
        w=w - 0.16,
        h=0.22,
        text=title,
        font_name=BODY_FONT,
        font_size=title_font_size,
        color=title_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    if embedded_label:
        add_textbox(
            slide,
            x=x + 0.08,
            y=y + h - 0.32,
            w=w - 0.16,
            h=0.18,
            text=embedded_label,
            font_name=FORMULA_FONT,
            font_size=embedded_label_font_size,
            color=embedded_label_color,
            bold=False,
            align=PP_ALIGN.CENTER,
        )


__all__ = [
    "add_box_title",
    "add_footer",
    "add_hub_box",
    "add_node_box",
    "add_rounded_box",
    "add_surface_card",
    "add_title_panel",
]
