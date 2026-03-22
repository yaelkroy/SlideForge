from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


def add_soft_connector(
    slide,
    *,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor,
    width_pt: float = 1.5,
) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line = connector.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    line.transparency = 0.15


def add_arrow(
    slide,
    *,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor,
    width_pt: float = 1.5,
) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line = connector.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    line.end_arrowhead = MSO_SHAPE.CHEVRON


__all__ = ["add_arrow", "add_soft_connector"]
