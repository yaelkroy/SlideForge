from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from slideforge.config.constants import (
    ACCENT,
    BODY_FONT,
    BOX_LINE,
    GHOST_TEXT,
    LIGHT_BOX_FILL,
    MUTED,
    NAVY,
    OFFWHITE,
    SLATE,
    TITLE_PANEL_FILL,
    TITLE_PANEL_LINE,
)
from slideforge.io.backgrounds import ensure_background_exists
from slideforge.utils.units import inches, pt


# -----------------------------------------------------------------------------
# Internal helpers
# -----------------------------------------------------------------------------


def _configure_text_frame(
    text_frame,
    *,
    word_wrap: bool = True,
    vertical_anchor=MSO_ANCHOR.TOP,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = word_wrap
    text_frame.vertical_anchor = vertical_anchor


def _apply_run_style(
    run,
    *,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    run.font.name = font_name
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_single_paragraph_text(
    text_frame,
    *,
    text: str,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    space_after_pt: int | None = None,
    line_spacing: float | None = None,
    first: bool = True,
):
    paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
    paragraph.alignment = align
    if space_after_pt is not None:
        paragraph.space_after = pt(space_after_pt)
    if line_spacing is not None:
        paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    _apply_run_style(
        run,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
    )
    return paragraph


def _add_text_to_shape(
    shape,
    *,
    text: str,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    tf = shape.text_frame
    _configure_text_frame(tf, vertical_anchor=vertical_anchor)
    _add_single_paragraph_text(
        tf,
        text=text,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
        first=True,
    )
    return tf


def _set_shape_line_and_fill(
    shape,
    *,
    line_color: RGBColor,
    fill_color: RGBColor | None,
    line_width_pt: float,
) -> None:
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = pt(line_width_pt)


def _add_shape(
    slide,
    *,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    line_color: RGBColor,
    fill_color: RGBColor | None,
    line_width_pt: float,
):
    shape = slide.shapes.add_shape(shape_type, inches(x), inches(y), inches(w), inches(h))
    _set_shape_line_and_fill(
        shape,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )
    return shape


# -----------------------------------------------------------------------------
# Backgrounds and text
# -----------------------------------------------------------------------------


def add_background(slide, prs, filename: str) -> None:
    bg_path = ensure_background_exists(filename)
    slide.shapes.add_picture(
        str(bg_path),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_name: str = BODY_FONT,
    font_size: int = 18,
    color: RGBColor = NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    _configure_text_frame(tf, vertical_anchor=vertical_anchor)
    _add_single_paragraph_text(
        tf,
        text=text,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
        first=True,
    )
    return box, tf


def style_paragraph(
    paragraph,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    for run in paragraph.runs:
        _apply_run_style(
            run,
            font_name=font_name,
            font_size=font_size,
            color=color,
            bold=bold,
        )


def add_bullets_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: list[str | tuple[str, int]],
    color: RGBColor = NAVY,
    top_font_size: int = 20,
    sub_font_size: int = 17,
    line_spacing: float = 1.15,
    top_space_after_pt: int = 8,
    sub_space_after_pt: int = 4,
) -> None:
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    _configure_text_frame(tf, vertical_anchor=MSO_ANCHOR.TOP)

    first = True
    for item in bullets:
        text, level = item if isinstance(item, tuple) else (item, 0)
        prefix = "• " if level == 0 else "– "
        indent = "    " * level
        paragraph = _add_single_paragraph_text(
            tf,
            text=f"{indent}{prefix}{text}",
            font_name=BODY_FONT,
            font_size=top_font_size if level == 0 else sub_font_size,
            color=color,
            bold=False,
            align=PP_ALIGN.LEFT,
            space_after_pt=top_space_after_pt if level == 0 else sub_space_after_pt,
            line_spacing=line_spacing,
            first=first,
        )
        style_paragraph(
            paragraph,
            font_name=BODY_FONT,
            font_size=top_font_size if level == 0 else sub_font_size,
            color=color,
            bold=False,
        )
        first = False


# -----------------------------------------------------------------------------
# Simple lines and cards
# -----------------------------------------------------------------------------


def add_divider_line(
    slide,
    dark: bool = False,
    *,
    x: float = 0.8,
    y: float = 1.1,
    w: float = 2.4,
    h: float = 0.04,
    color: RGBColor | None = None,
) -> None:
    actual_color = color or (OFFWHITE if dark else ACCENT)
    _add_shape(
        slide,
        shape_type=MSO_SHAPE.RECTANGLE,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=actual_color,
        fill_color=actual_color,
        line_width_pt=0.75,
    )


def add_header_rule(
    slide,
    *,
    x: float = 0.8,
    y: float = 1.1,
    w: float = 2.4,
    h: float = 0.04,
    color: RGBColor = ACCENT,
) -> None:
    add_divider_line(slide, x=x, y=y, w=w, h=h, color=color)


def add_rounded_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    line_color: RGBColor = BOX_LINE,
    fill_color: RGBColor | None = LIGHT_BOX_FILL,
    line_width_pt: float = 1.25,
):
    return _add_shape(
        slide,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )


def add_surface_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    line_color: RGBColor = BOX_LINE,
    fill_color: RGBColor | None = LIGHT_BOX_FILL,
    line_width_pt: float = 1.25,
):
    return add_rounded_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )


def add_box_title(
    slide,
    x: float,
    y: float,
    w: float,
    text: str,
    dark: bool = False,
    *,
    color: RGBColor | None = None,
    font_size: int = 11,
    bold: bool = True,
    align=PP_ALIGN.LEFT,
) -> None:
    actual_color = color or (OFFWHITE if dark else SLATE)
    add_textbox(
        slide,
        x=x,
        y=y,
        w=w,
        h=0.20,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=actual_color,
        bold=bold,
        align=align,
    )


def add_footer(
    slide,
    text: str = "Dr. Yael Demedetskaya",
    dark: bool = False,
    *,
    x: float = 0.45,
    y: float = 6.95,
    w: float = 4.2,
    h: float = 0.25,
    color: RGBColor | None = None,
    font_size: int = 10,
    align=PP_ALIGN.LEFT,
) -> None:
    actual_color = color or (OFFWHITE if dark else MUTED)
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    _add_text_to_shape(
        box,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=actual_color,
        bold=False,
        align=align,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


# -----------------------------------------------------------------------------
# Nodes, connectors, and labels
# -----------------------------------------------------------------------------


def add_node_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill_color: RGBColor = LIGHT_BOX_FILL,
    line_color: RGBColor = BOX_LINE,
    text_color: RGBColor = NAVY,
    font_size: int = 18,
    bold: bool = True,
    line_width_pt: float = 1.25,
):
    shape = add_rounded_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=line_width_pt,
    )
    _add_text_to_shape(
        shape,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=text_color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    return shape


def add_hub_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    fill_color: RGBColor = RGBColor(232, 238, 252),
    line_color: RGBColor = ACCENT,
    text_color: RGBColor = NAVY,
    font_size: int = 20,
    bold: bool = True,
):
    return add_node_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        text=text,
        fill_color=fill_color,
        line_color=line_color,
        text_color=text_color,
        font_size=font_size,
        bold=bold,
    )


def _add_straight_connector(
    slide,
    *,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor,
    width_pt: float,
    end_arrowhead: bool = False,
):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inches(x1),
        inches(y1),
        inches(x2),
        inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = pt(width_pt)
    if end_arrowhead:
        try:
            conn.line.end_arrowhead = True
        except Exception:
            pass
    return conn


def add_soft_connector(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor = TITLE_PANEL_LINE,
    width_pt: float = 1.5,
) -> None:
    _add_straight_connector(
        slide,
        x1=x1,
        y1=y1,
        x2=x2,
        y2=y2,
        color=color,
        width_pt=width_pt,
        end_arrowhead=False,
    )


def add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor = ACCENT,
    width_pt: float = 1.75,
):
    return _add_straight_connector(
        slide,
        x1=x1,
        y1=y1,
        x2=x2,
        y2=y2,
        color=color,
        width_pt=width_pt,
        end_arrowhead=True,
    )


def add_callout(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_size: int = 11,
    color: RGBColor = SLATE,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
):
    add_textbox(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
    )


def add_pill_tag(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill_color: RGBColor = TITLE_PANEL_FILL,
    line_color: RGBColor = TITLE_PANEL_LINE,
    text_color: RGBColor = OFFWHITE,
    font_size: int = 12,
    bold: bool = False,
):
    shape = add_rounded_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=1.0,
    )
    _add_text_to_shape(
        shape,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=text_color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_ghost_label(
    slide,
    x: float,
    y: float,
    w: float,
    text: str,
    font_size: int = 10,
    color: RGBColor = GHOST_TEXT,
) -> None:
    add_textbox(
        slide,
        x=x,
        y=y,
        w=w,
        h=0.2,
        text=text,
        font_name=BODY_FONT,
        font_size=font_size,
        color=color,
        bold=False,
        align=PP_ALIGN.CENTER,
    )


def add_title_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    embedded_label: str = "",
    *,
    fill_color: RGBColor = TITLE_PANEL_FILL,
    line_color: RGBColor = TITLE_PANEL_LINE,
    title_color: RGBColor = OFFWHITE,
    embedded_label_color: RGBColor = OFFWHITE,
    title_font_size: int = 11,
    embedded_label_font_size: int = 10,
) -> None:
    shape = add_rounded_box(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        line_color=line_color,
        fill_color=fill_color,
        line_width_pt=1.2,
    )

    add_textbox(
        slide,
        x=x + 0.08,
        y=y + 0.08,
        w=w - 0.16,
        h=0.22,
        text=title,
        font_name=BODY_FONT,
        font_size=title_font_size,
        color=title_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    if embedded_label:
        from slideforge.config.constants import FORMULA_FONT

        add_textbox(
            slide,
            x=x + 0.08,
            y=y + h - 0.32,
            w=w - 0.16,
            h=0.18,
            text=embedded_label,
            font_name=FORMULA_FONT,
            font_size=embedded_label_font_size,
            color=embedded_label_color,
            bold=False,
            align=PP_ALIGN.CENTER,
        )


__all__ = [
    "add_arrow",
    "add_background",
    "add_box_title",
    "add_bullets_box",
    "add_callout",
    "add_divider_line",
    "add_footer",
    "add_ghost_label",
    "add_header_rule",
    "add_hub_box",
    "add_node_box",
    "add_pill_tag",
    "add_rounded_box",
    "add_soft_connector",
    "add_surface_card",
    "add_textbox",
    "add_title_panel",
    "style_paragraph",
]
