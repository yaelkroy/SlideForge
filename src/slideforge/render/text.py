from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from slideforge.config.constants import BODY_FONT, NAVY
from slideforge.io.backgrounds import ensure_background_exists
from slideforge.utils.units import inches, pt


_ALIGN_LOOKUP = {
    "left": PP_ALIGN.LEFT,
    "l": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "centre": PP_ALIGN.CENTER,
    "c": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "r": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
    "justified": PP_ALIGN.JUSTIFY,
    "distributed": PP_ALIGN.DISTRIBUTE,
    "distribute": PP_ALIGN.DISTRIBUTE,
}

_VERTICAL_ANCHOR_LOOKUP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "center": MSO_ANCHOR.MIDDLE,
    "centre": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def _normalize_align(align):
    if align is None:
        return PP_ALIGN.LEFT
    if isinstance(align, str):
        return _ALIGN_LOOKUP.get(align.strip().lower(), PP_ALIGN.LEFT)
    try:
        return PP_ALIGN(align)
    except Exception:
        return PP_ALIGN.LEFT



def _normalize_vertical_anchor(vertical_anchor):
    if vertical_anchor is None:
        return MSO_ANCHOR.TOP
    if isinstance(vertical_anchor, str):
        return _VERTICAL_ANCHOR_LOOKUP.get(vertical_anchor.strip().lower(), MSO_ANCHOR.TOP)
    try:
        return MSO_ANCHOR(vertical_anchor)
    except Exception:
        return MSO_ANCHOR.TOP



def _configure_text_frame(
    text_frame,
    *,
    word_wrap: bool = True,
    vertical_anchor=MSO_ANCHOR.TOP,
) -> None:
    text_frame.clear()
    text_frame.word_wrap = word_wrap
    text_frame.vertical_anchor = _normalize_vertical_anchor(vertical_anchor)



def _apply_run_style(
    run,
    *,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    run.font.name = font_name
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color



def _add_single_paragraph_text(
    text_frame,
    *,
    text: str,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    space_after_pt: int | None = None,
    line_spacing: float | None = None,
    first: bool = True,
):
    paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
    paragraph.alignment = _normalize_align(align)
    if space_after_pt is not None:
        paragraph.space_after = pt(space_after_pt)
    if line_spacing is not None:
        paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    _apply_run_style(
        run,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
    )
    return paragraph



def _add_text_to_shape(
    shape,
    *,
    text: str,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    tf = shape.text_frame
    _configure_text_frame(tf, vertical_anchor=vertical_anchor)
    _add_single_paragraph_text(
        tf,
        text=text,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
        first=True,
    )
    return tf



def add_background(slide, prs, filename: str) -> None:
    bg_path = ensure_background_exists(filename)
    slide.shapes.add_picture(
        str(bg_path),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )



def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_name: str = BODY_FONT,
    font_size: int = 18,
    color: RGBColor = NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    _configure_text_frame(tf, vertical_anchor=vertical_anchor)
    _add_single_paragraph_text(
        tf,
        text=text,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
        first=True,
    )
    return box, tf



def style_paragraph(
    paragraph,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    for run in paragraph.runs:
        _apply_run_style(
            run,
            font_name=font_name,
            font_size=font_size,
            color=color,
            bold=bold,
        )



def add_bullets_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    bullets: list[str | tuple[str, int]],
    color: RGBColor = NAVY,
    top_font_size: int = 20,
    sub_font_size: int = 17,
    line_spacing: float = 1.15,
    top_space_after_pt: int = 8,
    sub_space_after_pt: int = 4,
) -> None:
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    _configure_text_frame(tf, vertical_anchor=MSO_ANCHOR.TOP)

    first = True
    for item in bullets:
        text, level = item if isinstance(item, tuple) else (item, 0)
        prefix = "• " if level == 0 else "– "
        indent = "    " * level
        paragraph = _add_single_paragraph_text(
            tf,
            text=f"{indent}{prefix}{text}",
            font_name=BODY_FONT,
            font_size=top_font_size if level == 0 else sub_font_size,
            color=color,
            bold=False,
            align=PP_ALIGN.LEFT,
            space_after_pt=top_space_after_pt if level == 0 else sub_space_after_pt,
            line_spacing=line_spacing,
            first=first,
        )
        style_paragraph(
            paragraph,
            font_name=BODY_FONT,
            font_size=top_font_size if level == 0 else sub_font_size,
            color=color,
            bold=False,
        )
        first = False


__all__ = [
    "add_background",
    "add_bullets_box",
    "add_textbox",
    "style_paragraph",
]
