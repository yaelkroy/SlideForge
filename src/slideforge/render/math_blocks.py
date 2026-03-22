from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Mapping, Sequence

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from slideforge.config.constants import (
    BODY_FONT,
    FORMULA_FONT,
    LIGHT_BOX_FILL,
    NAVY,
    OFFWHITE,
    SLATE,
)
from slideforge.layout.autofit import Box, estimate_text_height, fit_text
from slideforge.render.primitives import add_box_title, add_rounded_box, style_paragraph


@dataclass(frozen=True)
class MathBlockRenderResult:
    box: Box
    font_size: int
    line_count: int
    text: str


@dataclass(frozen=True)
class MathBlockStyle:
    body_font_name: str = BODY_FONT
    formula_font_name: str = FORMULA_FONT
    body_color: RGBColor = SLATE
    formula_color: RGBColor = NAVY
    result_color: RGBColor = NAVY
    label_color: RGBColor = SLATE
    card_fill_color: RGBColor | None = LIGHT_BOX_FILL
    card_line_color: RGBColor = RGBColor(173, 185, 220)
    final_answer_fill_color: RGBColor | None = OFFWHITE


# Conservative safety margins used to avoid the last line touching or escaping the box.
_BOTTOM_SAFETY_IN = 0.055
_TOP_SAFETY_IN = 0.02
_SIDE_SAFETY_IN = 0.01
_DEFAULT_TEXT_SAFETY_RATIO = 0.89
_RESULT_TEXT_SAFETY_RATIO = 0.86


def _clean_text(value: Any) -> str:
    return str(value or "").strip()


def _clean_lines(items: Sequence[Any] | None) -> list[str]:
    return [_clean_text(item) for item in (items or []) if _clean_text(item)]


def _normalize_formula_lines(formulas: str | Sequence[Any]) -> list[str]:
    if isinstance(formulas, str):
        return [line.strip() for line in formulas.splitlines() if line.strip()]
    return _clean_lines(formulas)


def _normalize_step(step: Any, index: int) -> dict[str, str]:
    if isinstance(step, Mapping):
        return {
            "title": _clean_text(step.get("title") or step.get("label") or f"Step {index}"),
            "body": _clean_text(step.get("body") or step.get("text") or step.get("explanation")),
            "formula": _clean_text(step.get("formula") or step.get("equation") or step.get("result")),
            "note": _clean_text(step.get("note")),
        }
    return {
        "title": f"Step {index}",
        "body": _clean_text(step),
        "formula": "",
        "note": "",
    }


def _normalize_steps(steps: Sequence[Any] | None) -> list[dict[str, str]]:
    return [_normalize_step(step, idx + 1) for idx, step in enumerate(steps or [])]


def _shrink_box(
    box: Box,
    *,
    top: float = _TOP_SAFETY_IN,
    bottom: float = _BOTTOM_SAFETY_IN,
    side: float = _SIDE_SAFETY_IN,
    ratio: float = _DEFAULT_TEXT_SAFETY_RATIO,
) -> Box:
    if box.w <= 0 or box.h <= 0:
        return box

    inner_x = box.x + side
    inner_y = box.y + top
    inner_w = max(0.0, box.w - 2 * side)
    inner_h = max(0.0, box.h - top - bottom)
    safe_h = max(0.0, inner_h * ratio)
    return Box(inner_x, inner_y, inner_w, safe_h)


def _estimate_safe_height(
    text: str,
    *,
    width: float,
    font_size: int,
    line_spacing: float,
    min_lines: int = 1,
) -> float:
    if not text or width <= 0:
        return 0.0
    estimated = estimate_text_height(
        text,
        width,
        font_size=font_size,
        line_spacing=line_spacing,
    )
    min_height = min_lines * (font_size / 72.0) * line_spacing
    return max(estimated, min_height) + _BOTTOM_SAFETY_IN


def _fit_font_size(
    text: str,
    box: Box,
    *,
    min_font: int,
    max_font: int,
    max_lines: int | None,
    line_spacing: float = 1.12,
    prefer_single_line: bool = False,
    safety_ratio: float = _DEFAULT_TEXT_SAFETY_RATIO,
    extra_shrink_loops: int = 2,
) -> int:
    if not text.strip() or box.w <= 0 or box.h <= 0:
        return max_font

    safe_box = _shrink_box(box, ratio=safety_ratio)
    fitted = fit_text(
        text,
        safe_box.w,
        safe_box.h,
        min_font_size=min_font,
        max_font_size=max_font,
        max_lines=max_lines,
        line_spacing=line_spacing,
        prefer_single_line=prefer_single_line,
    )
    font_size = max(min_font, fitted.font_size)

    # Extra conservative shrink loop in case the initial fit is too optimistic.
    for _ in range(max(0, extra_shrink_loops)):
        est_h = _estimate_safe_height(
            text,
            width=safe_box.w,
            font_size=font_size,
            line_spacing=line_spacing,
            min_lines=1,
        )
        if est_h <= safe_box.h + 1e-6:
            break
        font_size -= 1
        if font_size <= min_font:
            font_size = min_font
            break

    return max(min_font, font_size)


def _add_text_frame(slide, *, box: Box, vertical_anchor=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    return shape, tf


def _add_paragraph(
    tf,
    *,
    text: str,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    space_after_pt: int = 4,
    line_spacing: float = 1.12,
    first: bool = False,
):
    paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    paragraph.space_after = Pt(space_after_pt)
    run = paragraph.add_run()
    run.text = text
    style_paragraph(
        paragraph,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
    )
    return paragraph


def estimate_multiline_formula_height(
    formulas: str | Sequence[Any],
    *,
    width: float,
    font_size: int = 16,
    line_spacing: float = 1.10,
) -> float:
    lines = _normalize_formula_lines(formulas)
    if not lines:
        return 0.0
    return sum(
        _estimate_safe_height(line, width=width, font_size=font_size, line_spacing=line_spacing)
        for line in lines
    )


def estimate_derivation_height(
    steps: Sequence[Any],
    *,
    width: float,
    body_font_size: int = 12,
    formula_font_size: int = 13,
) -> float:
    normalized_steps = _normalize_steps(steps)
    if not normalized_steps:
        return 0.0
    total = 0.0
    for step in normalized_steps:
        total += _estimate_safe_height(step["title"], width=width, font_size=body_font_size, line_spacing=1.08)
        if step["body"]:
            total += _estimate_safe_height(step["body"], width=width, font_size=body_font_size, line_spacing=1.08)
        if step["formula"]:
            total += _estimate_safe_height(step["formula"], width=width, font_size=formula_font_size, line_spacing=1.06)
        if step["note"]:
            total += _estimate_safe_height(step["note"], width=width, font_size=max(10, body_font_size - 1), line_spacing=1.06)
    return total


def estimate_result_callout_height(
    result_lines: str | Sequence[Any],
    *,
    width: float,
    font_size: int = 14,
    label_h: float = 0.22,
) -> float:
    lines = _normalize_formula_lines(result_lines)
    if not lines:
        return 0.0
    body_h = sum(
        _estimate_safe_height(line, width=width, font_size=font_size + (1 if idx == len(lines) - 1 else 0), line_spacing=1.08)
        for idx, line in enumerate(lines)
    )
    return label_h + body_h + 0.16


def render_multiline_formulas(
    slide,
    *,
    box: Box,
    formulas: str | Sequence[Any],
    style: MathBlockStyle | None = None,
    min_font: int = 14,
    max_font: int = 22,
    max_lines: int | None = 8,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.10,
    emphasize_last: bool = False,
) -> MathBlockRenderResult:
    """Render one or more formulas as native PowerPoint text with safer bottom padding."""
    style = style or MathBlockStyle()
    lines = _normalize_formula_lines(formulas)
    text = "\n".join(lines)
    safe_box = _shrink_box(box)
    font_size = _fit_font_size(
        text,
        safe_box,
        min_font=min_font,
        max_font=max_font,
        max_lines=max_lines,
        line_spacing=line_spacing,
        safety_ratio=_DEFAULT_TEXT_SAFETY_RATIO,
    )
    _, tf = _add_text_frame(slide, box=safe_box)
    for idx, line in enumerate(lines):
        is_last = idx == len(lines) - 1
        _add_paragraph(
            tf,
            text=line,
            font_name=style.formula_font_name,
            font_size=font_size + (1 if emphasize_last and is_last else 0),
            color=style.formula_color,
            bold=bool(emphasize_last and is_last),
            align=align,
            space_after_pt=3 if not is_last else 0,
            line_spacing=line_spacing,
            first=idx == 0,
        )
    return MathBlockRenderResult(box=box, font_size=font_size, line_count=len(lines), text=text)


def render_compact_derivation_stack(
    slide,
    *,
    box: Box,
    steps: Sequence[Any],
    style: MathBlockStyle | None = None,
    min_body_font: int = 11,
    max_body_font: int = 15,
    min_formula_font: int = 12,
    max_formula_font: int = 16,
    final_answer: str = "",
    emphasize_final_answer: bool = True,
    final_answer_style: str = "inline",
    align=PP_ALIGN.LEFT,
) -> MathBlockRenderResult:
    """Render a compact step-by-step derivation stack with conservative fit rules."""
    style = style or MathBlockStyle()
    normalized_steps = _normalize_steps(steps)

    plain_lines: list[str] = []
    for step in normalized_steps:
        plain_lines.append(step["title"])
        if step["body"]:
            plain_lines.append(step["body"])
        if step["formula"]:
            plain_lines.append(step["formula"])
        if step["note"]:
            plain_lines.append(step["note"])
    final_answer_text = _clean_text(final_answer)
    if final_answer_text and final_answer_style == "inline":
        plain_lines.append(final_answer_text)

    text = "\n".join(plain_lines)
    safe_box = _shrink_box(box)
    body_font = _fit_font_size(
        text,
        safe_box,
        min_font=min_body_font,
        max_font=max_body_font,
        max_lines=max(8, len(plain_lines) + 2),
        line_spacing=1.08,
        safety_ratio=_DEFAULT_TEXT_SAFETY_RATIO,
    )
    formula_font = max(min_formula_font, min(max_formula_font, body_font + 1))

    if final_answer_text and final_answer_style == "callout":
        callout_h = min(max(0.40, safe_box.h * 0.18), 0.75)
        steps_box = Box(safe_box.x, safe_box.y, safe_box.w, max(0.0, safe_box.h - callout_h - 0.06))
        callout_box = Box(safe_box.x, steps_box.y + steps_box.h + 0.06, safe_box.w, callout_h)
    else:
        steps_box = safe_box
        callout_box = None

    _, tf = _add_text_frame(slide, box=steps_box)
    first = True
    for step in normalized_steps:
        _add_paragraph(
            tf,
            text=step["title"],
            font_name=style.body_font_name,
            font_size=body_font,
            color=style.label_color,
            bold=True,
            align=align,
            space_after_pt=2,
            line_spacing=1.08,
            first=first,
        )
        first = False
        if step["body"]:
            _add_paragraph(
                tf,
                text=step["body"],
                font_name=style.body_font_name,
                font_size=body_font,
                color=style.body_color,
                align=align,
                space_after_pt=2,
                line_spacing=1.08,
            )
        if step["formula"]:
            _add_paragraph(
                tf,
                text=step["formula"],
                font_name=style.formula_font_name,
                font_size=formula_font,
                color=style.formula_color,
                align=align,
                space_after_pt=2,
                line_spacing=1.06,
            )
        if step["note"]:
            _add_paragraph(
                tf,
                text=step["note"],
                font_name=style.body_font_name,
                font_size=max(min_body_font, body_font - 1),
                color=style.body_color,
                align=align,
                space_after_pt=4,
                line_spacing=1.05,
            )

    if final_answer_text:
        if callout_box is not None:
            render_result_callout(
                slide,
                box=callout_box,
                result_lines=[final_answer_text],
                label="Result",
                style=style,
                min_font=max(min_formula_font, formula_font),
                max_font=max_formula_font + 1,
                emphasize_final_answer=True,
                align=align,
                line_spacing=1.05,
                draw_card=True,
            )
        else:
            _add_paragraph(
                tf,
                text=final_answer_text,
                font_name=style.formula_font_name,
                font_size=max(formula_font, body_font + 1),
                color=style.result_color,
                bold=emphasize_final_answer,
                align=align,
                space_after_pt=0,
                line_spacing=1.05,
            )

    return MathBlockRenderResult(box=box, font_size=body_font, line_count=len(plain_lines), text=text)


def render_result_callout(
    slide,
    *,
    box: Box,
    result_lines: str | Sequence[Any],
    label: str = "Result",
    style: MathBlockStyle | None = None,
    min_font: int = 13,
    max_font: int = 18,
    emphasize_final_answer: bool = True,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.08,
    draw_card: bool = True,
) -> MathBlockRenderResult:
    """Render a result strip/callout with stronger final-answer protection."""
    style = style or MathBlockStyle()
    lines = _normalize_formula_lines(result_lines)
    if not lines:
        lines = [""]

    if draw_card:
        add_rounded_box(
            slide,
            box.x,
            box.y,
            box.w,
            box.h,
            line_color=style.card_line_color,
            fill_color=style.card_fill_color,
            line_width_pt=1.25,
        )

    label_box = Box(box.x + 0.16, box.y + 0.08, max(0.0, box.w - 0.32), 0.20)
    add_box_title(
        slide,
        x=label_box.x,
        y=label_box.y,
        w=label_box.w,
        text=label,
        color=style.label_color,
        font_size=11,
        bold=True,
        align=PP_ALIGN.LEFT,
    )

    body_box = Box(box.x + 0.18, box.y + 0.34, max(0.0, box.w - 0.36), max(0.0, box.h - 0.48))
    body_safe_box = _shrink_box(body_box, ratio=_RESULT_TEXT_SAFETY_RATIO, bottom=max(_BOTTOM_SAFETY_IN, 0.075))
    text = "\n".join(lines)
    font_size = _fit_font_size(
        text,
        body_safe_box,
        min_font=min_font,
        max_font=max_font,
        max_lines=max(3, len(lines) + 1),
        line_spacing=line_spacing,
        safety_ratio=_RESULT_TEXT_SAFETY_RATIO,
        extra_shrink_loops=3,
    )

    _, tf = _add_text_frame(slide, box=body_safe_box, vertical_anchor=MSO_ANCHOR.MIDDLE)
    for idx, line in enumerate(lines):
        is_last = idx == len(lines) - 1
        _add_paragraph(
            tf,
            text=line,
            font_name=style.formula_font_name,
            font_size=font_size + (1 if is_last and emphasize_final_answer else 0),
            color=style.result_color,
            bold=bool(is_last and emphasize_final_answer),
            align=align,
            space_after_pt=2 if not is_last else 0,
            line_spacing=1.05 if is_last else line_spacing,
            first=idx == 0,
        )

    return MathBlockRenderResult(box=box, font_size=font_size, line_count=len(lines), text=text)


def render_emphasized_final_answer(
    slide,
    *,
    box: Box,
    text: str,
    style: MathBlockStyle | None = None,
    min_font: int = 14,
    max_font: int = 20,
    align=PP_ALIGN.CENTER,
) -> MathBlockRenderResult:
    """Render a small highlighted answer pill/card with conservative fit margins."""
    style = style or MathBlockStyle()
    add_rounded_box(
        slide,
        box.x,
        box.y,
        box.w,
        box.h,
        line_color=style.card_line_color,
        fill_color=style.final_answer_fill_color,
        line_width_pt=1.0,
    )

    safe_box = _shrink_box(box, ratio=_RESULT_TEXT_SAFETY_RATIO, bottom=max(_BOTTOM_SAFETY_IN, 0.07))
    font_size = _fit_font_size(
        text,
        safe_box,
        min_font=min_font,
        max_font=max_font,
        max_lines=2,
        line_spacing=1.03,
        prefer_single_line=False,
        safety_ratio=_RESULT_TEXT_SAFETY_RATIO,
        extra_shrink_loops=3,
    )
    _, tf = _add_text_frame(slide, box=safe_box, vertical_anchor=MSO_ANCHOR.MIDDLE)
    _add_paragraph(
        tf,
        text=text,
        font_name=style.formula_font_name,
        font_size=font_size,
        color=style.result_color,
        bold=True,
        align=align,
        space_after_pt=0,
        line_spacing=1.03,
        first=True,
    )
    return MathBlockRenderResult(box=box, font_size=font_size, line_count=1, text=text)


__all__ = [
    "MathBlockRenderResult",
    "MathBlockStyle",
    "estimate_multiline_formula_height",
    "estimate_derivation_height",
    "estimate_result_callout_height",
    "render_multiline_formulas",
    "render_compact_derivation_stack",
    "render_result_callout",
    "render_emphasized_final_answer",
]
